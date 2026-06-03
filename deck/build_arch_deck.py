# Task Management System | Software Architecture deck (SEN3006).
# V3 redesign: warm editorial system. Spectral serif headers + Inter body, a
# warm cream base, an indigo spine with three pattern hues (Factory=clay,
# Strategy=teal, State=violet) plus status colours, structural Untitled UI icons
# on every card, colour-coded section dividers, and a section progress tracker.
# Presenter names appear on the cover only. Built with python-pptx.
# Build: python build_arch_deck.py   (icons live in _ico/, diagrams in _dia/)
import os, re
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
DIA = os.path.join(HERE, "_dia")
ICO = os.path.join(HERE, "_ico")
OUT = os.path.join(HERE, "Task-Management-System.pptx")

# ---- palette ----------------------------------------------------------------
def C(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
PAPER  = C("FBF9F5")   # content background, warm cream
CARD   = C("FFFFFF")   # card fill
CARDBD = C("E8E2D5")   # card border
RULE   = C("E4DECF")   # hairlines on cream
INK    = C("1C1A16")   # serif headers + dark backgrounds
BODY   = C("57534B")   # Inter body on cream
MUTE   = C("8C877D")   # labels, captions
FAINT  = C("B7B1A4")   # footers, de-emphasis
DARK   = C("1C1A16")   # divider / cover background
RULE_DK= C("3A352E")   # hairlines on dark
WHITE  = C("FFFFFF")
GRAYW  = C("A8A299")   # body text on dark
PANEL  = C("F4F0E8")   # subtle panel on cream

# accents on light (deeper)
INDIGO = C("4A45C7"); CLAY = C("BC5A2B"); TEAL = C("0E7C72"); VIOLET = C("7C3AED")
GREEN  = C("1E7A45"); AMBER = C("B5760C"); RED  = C("BB2D23")
# tints (chip backgrounds)
INDIGO_T=C("ECEBFA"); CLAY_T=C("F6ECE3"); TEAL_T=C("E2F0ED"); VIOLET_T=C("F0EAFB")
GREEN_T =C("E6F0EA"); AMBER_T=C("F5EEDD"); RED_T =C("F7E9E6"); MUTE_T=C("EEEBE6")
# accents on dark (brighter) for dividers
INDIGO_D=C("8E8BF5"); VIOLET_D=C("B6A6F7"); TEAL_D=C("4FC2B2"); GREEN_D=C("5BD08F"); AMBER_D=C("E7B765")

# code syntax (muted editor grade)
CODEBG=C("F7F4ED"); CODEBD=C("E7E1D4")
KW=C("8250DF"); TYP=C("0550AE"); STRc=C("0A7B33"); NUMc=C("0550AE"); COM=C("8C877D")

SERIF_SB="Spectral SemiBold"; SERIF="Spectral"; SERIF_MD="Spectral Medium"
SANS="Inter"; SANS_SB="Inter SemiBold"; SANS_MD="Inter Medium"; MONO="Consolas"

SW, SH = 13.333, 7.5; M = 0.82; TOTAL = 32

prs = Presentation(); prs.slide_width = Emu(int(SW*914400)); prs.slide_height = Emu(int(SH*914400))
BLANK = prs.slide_layouts[6]

# ---- primitives -------------------------------------------------------------
def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb = bg; return s
def notes(s, txt): s.notes_slide.notes_text_frame.text = txt
def _spc(r, pts): r.font._rPr.set("spc", str(int(pts*100)))
def hline(s, x, y, w, color=RULE, weight=1.0):
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x+w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(weight); ln.shadow.inherit = False; return ln
def rrect(s, x, y, w, h, fill, line=None, lw=1.0, radius=0.06):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    try: sp.adjustments[0] = radius
    except Exception: pass
    return sp
def text(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tf = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf, m, 0)
    first = True
    for para in paras:
        runs = para if isinstance(para, list) else [para]
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False; p.alignment = align; meta = runs[0]
        if meta.get("space_before") is not None: p.space_before = Pt(meta["space_before"])
        if meta.get("space_after") is not None: p.space_after = Pt(meta["space_after"])
        if meta.get("line") is not None: p.line_spacing = meta["line"]
        for rd in runs:
            r = p.add_run(); r.text = rd["t"]; f = r.font
            f.name = rd.get("font", SANS); f.size = Pt(rd.get("size", 13.5))
            f.bold = rd.get("bold", False); f.color.rgb = rd.get("color", INK)
            if rd.get("italic"): f.italic = True
            if rd.get("track"): _spc(r, rd["track"])
    return tf
def place_icon(s, name, key, cx, cy, size):
    p = os.path.join(ICO, f"{name}__{key}.png")
    if not os.path.exists(p):
        return
    s.shapes.add_picture(p, Inches(cx-size/2), Inches(cy-size/2), Inches(size), Inches(size))
def chip(s, x, y, size, tint, name, key, frac=0.56):
    rrect(s, x, y, size, size, tint, radius=0.30)
    place_icon(s, name, key, x+size/2, y+size/2, size*frac)

# ---- chrome -----------------------------------------------------------------
def eyebrow(s, label, accent=INDIGO):
    rrect(s, M, 0.64, 0.07, 0.24, accent, radius=0.5)
    text(s, M+0.22, 0.64, 9, 0.3, [{"t": label.upper(), "size": 11, "color": accent, "bold": True, "track": 2.0, "font": SANS_SB}])
SECTION_HUE = {1: INDIGO, 2: VIOLET, 3: TEAL, 4: GREEN, 5: AMBER}
def tracker(s, active, on_dark=False):
    n=5; segw=0.32; gap=0.10; y=0.70
    x0 = SW-M-(n*segw+(n-1)*gap)
    hue = SECTION_HUE.get(active, INDIGO)
    for i in range(n):
        x = x0+i*(segw+gap)
        if i == active-1: col = hue
        else:             col = RULE_DK if on_dark else RULE
        rrect(s, x, y, segw, 0.07, col, radius=0.5)
    text(s, x0, y+0.13, n*segw+(n-1)*gap, 0.2,
         [{"t": f"SECTION {active} / 5", "size": 7.5, "color": hue, "track": 1.4, "font": SANS_SB}],
         align=PP_ALIGN.RIGHT)
def title(s, t, y=1.42, color=INK, size=29, w=None):
    text(s, M, y, w or (SW-2*M), 1.0, [{"t": t, "size": size, "color": color, "font": SERIF_SB, "line": 1.04}])
def sub(s, t, color=BODY, y=2.26, w=None, size=13):
    text(s, M, y, w or (SW-2*M-1.0), 0.7, [{"t": t, "size": size, "color": color, "line": 1.36}])
def footer(s, n, on_dark=False):
    c = GRAYW if on_dark else FAINT
    text(s, M, SH-0.52, 8, 0.3, [{"t": "Task Management System   ·   SEN3006 Software Architecture", "size": 8.5, "color": c, "track": 1.0}])
    text(s, SW-M-3, SH-0.52, 3, 0.3, [{"t": f"{n:02d} / {TOTAL}", "size": 8.5, "color": c, "track": 1.2}], align=PP_ALIGN.RIGHT)

# ---- composite blocks -------------------------------------------------------
def card(s, x, y, w, h, tint, name, key, title_t, body_t, tsize=13.5, bsize=11.5, cs=0.5):
    rrect(s, x, y, w, h, CARD, CARDBD, 1.0, radius=0.045)
    chip(s, x+0.28, y+0.28, cs, tint, name, key)
    text(s, x+0.28+cs+0.18, y+0.28, w-(0.28+cs+0.18)-0.24, cs,
         [{"t": title_t, "size": tsize, "color": INK, "font": SANS_SB, "line": 1.04}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+0.30, y+0.28+cs+0.14, w-0.58, h-(0.28+cs+0.14)-0.16,
         [{"t": body_t, "size": bsize, "color": BODY, "line": 1.32}])
def listpanel(s, x, y, w, h, tint, name, key, label, items, acc, isize=12, gap=0.6):
    rrect(s, x, y, w, h, CARD, CARDBD, 1.0, radius=0.045)
    chip(s, x+0.30, y+0.28, 0.5, tint, name, key)
    text(s, x+0.30+0.5+0.18, y+0.28, w-1.2, 0.5,
         [{"t": label.upper(), "size": 11, "color": INK, "font": SANS_SB, "track": 1.4}], anchor=MSO_ANCHOR.MIDDLE)
    for i, it in enumerate(items):
        yy = y+1.04+i*gap
        text(s, x+0.34, yy-0.01, 0.22, 0.3, [{"t": "›", "size": 11.5, "color": acc, "font": SANS_SB}])
        text(s, x+0.58, yy, w-0.86, gap+0.1, [{"t": it, "size": isize, "color": BODY, "line": 1.2}])
def iconcol3(s, triples, y, specs, rule=True):
    gap=0.5; colw=(SW-2*M-2*gap)/3
    for i,(h,b) in enumerate(triples):
        x = M+i*(colw+gap); acc,tint,name,key = specs[i]
        chip(s, x, y, 0.5, tint, name, key)
        text(s, x, y+0.64, colw, 0.4, [{"t": h, "size": 13, "color": INK, "font": SANS_SB}])
        if rule: hline(s, x, y+0.97, colw*0.42, acc, 2.2)
        text(s, x, y+1.10, colw, 1.4, [{"t": b, "size": 11.5, "color": BODY, "line": 1.32}])
def console(s, x, y, w, h, rows, header="bash"):
    """Dark terminal card: window dots, header, then monospace rows of (text, color)."""
    rrect(s, x, y, w, h, C("15140F"), radius=0.05)
    for i,c in enumerate(["E66A5A","E7B765","5BD08F"]):
        rrect(s, x+0.22+i*0.20, y+0.135, 0.085, 0.085, C(c), radius=0.5)
    text(s, x, y+0.05, w, 0.26, [{"t": header, "size": 8.5, "color": C("6E6A60"), "font": MONO}], align=PP_ALIGN.CENTER)
    hline(s, x, y+0.36, w, C("2C291F"), 1.0)
    tf = s.shapes.add_textbox(Inches(x+0.24), Inches(y+0.46), Inches(w-0.42), Inches(h-0.56)).text_frame
    tf.word_wrap=False; tf.vertical_anchor=MSO_ANCHOR.TOP
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf, m, 0)
    first=True
    for txt,col in rows:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.line_spacing=1.24; p.alignment=PP_ALIGN.LEFT
        r=p.add_run(); r.text=txt if txt else " "; r.font.name=MONO; r.font.size=Pt(9); r.font.color.rgb=col

def cmdcard(s, x, y, w, label, command, accent, csize=13, h=0.92):
    """A 'paste this' command block: accent label on top, the exact command in mono below."""
    rrect(s, x, y, w, h, CODEBG, CODEBD, 1.0, radius=0.06)
    hline(s, x, y, w, accent, 2.4)
    text(s, x+0.26, y+0.16, w-0.5, 0.3, [{"t": label, "size": 10, "color": accent, "bold": True, "track": 1.5, "font": SANS_SB}])
    text(s, x+0.26, y+h-0.44, w-0.5, 0.36, [{"t": command, "size": csize, "color": INK, "font": MONO}])

def img_panel(s, name, bx, by, bw, bh, accent=None, caption=None):
    """White panel behind a diagram/screenshot, optional accent top-rule + caption."""
    rrect(s, bx, by, bw, bh, CARD, CARDBD, 1.0, radius=0.03)
    if accent: hline(s, bx+0.0, by, bw, accent, 2.4)
    pad = 0.18; ix, iy, iw, ih = bx+pad, by+pad+0.04, bw-2*pad, bh-2*pad-0.04
    path = os.path.join(DIA, name)
    if not os.path.exists(path):
        text(s, ix, iy+ih/2-0.2, iw, 0.4, [{"t": "[ "+name+" ]", "size": 12, "color": FAINT}], align=PP_ALIGN.CENTER); return
    iwpx, ihpx = Image.open(path).size; ratio = iwpx/ihpx
    if iw/ih > ratio: h = ih; w = ih*ratio
    else: w = iw; h = iw/ratio
    s.shapes.add_picture(path, Inches(ix+(iw-w)/2), Inches(iy+(ih-h)/2), Inches(w), Inches(h))
    if caption:
        text(s, bx, by+bh+0.1, bw, 0.3, [{"t": caption, "size": 10, "color": MUTE, "track": 0.4}], align=PP_ALIGN.CENTER)

def add_arrow(conn):
    ln = conn.line._get_or_add_ln(); t = ln.find(qn("a:tailEnd"))
    if t is None: t = ln.makeelement(qn("a:tailEnd"), {}); ln.append(t)
    t.set("type","triangle"); t.set("w","med"); t.set("len","med")

def blendhex(h1, h2, t):
    a=[int(h1[i:i+2],16) for i in (0,2,4)]; b=[int(h2[i:i+2],16) for i in (0,2,4)]
    return RGBColor(*[int(round(a[i]+(b[i]-a[i])*t)) for i in range(3)])
def divider(s, kicker, statement, sub_text, hexc, icon_name, active):
    # full-bleed section colour; white icon and statement, dimmed-white supporting marks
    s.background.fill.solid(); s.background.fill.fore_color.rgb = C(hexc)
    dim  = blendhex(hexc, "FFFFFF", 0.82)   # secondary text on the colour field
    rule = blendhex(hexc, "FFFFFF", 0.48)   # hairline on the colour field
    soff = blendhex(hexc, "FFFFFF", 0.30)   # inactive tracker segments
    text(s, M, 0.66, 9, 0.3, [{"t": kicker.upper(), "size": 11, "color": WHITE, "bold": True, "track": 2.8, "font": SANS_SB}])
    n=5; segw=0.32; gap=0.10; ty=0.70; x0=SW-M-(n*segw+(n-1)*gap)
    for i in range(n):
        rrect(s, x0+i*(segw+gap), ty, segw, 0.07, WHITE if i==active-1 else soff, radius=0.5)
    text(s, x0, ty+0.13, n*segw+(n-1)*gap, 0.2,
         [{"t": f"SECTION {active} / 5", "size": 7.5, "color": dim, "track": 1.4, "font": SANS_SB}], align=PP_ALIGN.RIGHT)
    place_icon(s, icon_name, "white", SW-M-1.25, 2.05, 1.55)
    hline(s, M, 3.30, SW-2*M, rule, 1.6)
    text(s, M, 3.60, SW-2*M-1.4, 1.4, [{"t": statement, "size": 37, "color": WHITE, "font": SERIF_SB, "line": 1.05}])
    text(s, M, 4.78, SW-2*M-2.2, 1.0, [{"t": sub_text, "size": 14, "color": dim, "line": 1.42}])

# ---- Java syntax highlighter + code panel -----------------------------------
JAVA_KW = {"class","interface","enum","abstract","extends","implements","new","void","return",
           "switch","case","int","boolean","public","private","protected","static","final",
           "this","throws","import","package","if","else"}
_tok = re.compile(r'(//[^\n]*|"[^"]*"|[A-Za-z_][A-Za-z0-9_]*|\d+|\s+|.)')
def java_runs(line):
    if not line.strip(): return [{"t": " ", "color": INK}]
    out=[]
    for m in _tok.finditer(line):
        t=m.group(0)
        if t.startswith("//"): col=COM
        elif t.startswith('"'): col=STRc
        elif t.isspace(): col=INK
        elif t.isdigit(): col=NUMc
        elif t in JAVA_KW: col=KW
        elif re.match(r'^[A-Z][A-Za-z0-9_]*$', t): col=TYP
        else: col=INK
        out.append({"t": t, "color": col})
    return out
def code(s, x, y, w, lines, size=12.5, accent=INDIGO):
    lh=size/72*1.55; h=0.40+len(lines)*lh
    rrect(s, x, y, w, h, CODEBG, CODEBD, 1.0, radius=0.04)
    hline(s, x, y, w, accent, 2.2)
    tf = s.shapes.add_textbox(Inches(x+0.04), Inches(y+0.14), Inches(w-0.08), Inches(h-0.2)).text_frame
    tf.word_wrap=False; tf.vertical_anchor=MSO_ANCHOR.TOP
    tf.margin_left=Inches(0.26); tf.margin_right=Inches(0.18); tf.margin_top=Inches(0.06); tf.margin_bottom=Inches(0.06)
    first=True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.line_spacing=1.42; p.alignment=PP_ALIGN.LEFT
        for rd in java_runs(ln):
            r=p.add_run(); r.text=rd["t"]; r.font.name=MONO; r.font.size=Pt(size); r.font.color.rgb=rd["color"]
    return h

# ============================== SLIDES ==============================
NAME1="Abdul Rahman Malak"; NAME2="Mouhammad Houjeirat"

# 01 COVER (names live here only) --------------------------------------------
s = slide(DARK)
text(s, M, 0.66, 9, 0.3, [{"t": "SEN3006  ·  SOFTWARE ARCHITECTURE", "size": 11, "color": GRAYW, "track": 2.4, "font": SANS_SB}])
text(s, SW-M-6, 0.66, 6, 0.3, [{"t": "PROJECT PRESENTATION", "size": 11, "color": GRAYW, "track": 1.8, "font": SANS_SB}], align=PP_ALIGN.RIGHT)
text(s, M, 2.35, 11.9, 1.5, [{"t": "Task Management System.", "size": 50, "color": WHITE, "font": SERIF_SB, "line": 1.0}])
text(s, M, 3.62, 11.9, 0.7, [{"t": "Factory Method and Strategy, built in pure Java.", "size": 21, "color": INDIGO_D, "font": SERIF_MD, "italic": True, "line": 1.2}])
# pattern colour key rail
keyrow = [("Factory Method", CLAY), ("Strategy", TEAL), ("State", VIOLET)]
kx = M
for lab, col in keyrow:
    rrect(s, kx, 4.78, 0.30, 0.10, col, radius=0.5)
    text(s, kx+0.40, 4.66, 2.6, 0.34, [{"t": lab, "size": 11.5, "color": GRAYW, "font": SANS_SB, "track": 0.6}])
    kx += 0.40 + (1.95 if lab!="Strategy" else 1.45)
hline(s, M, 5.92, SW-2*M, RULE_DK, 1.0)
text(s, M, 6.14, 7.6, 0.8, [{"t": "A task tracker for software teams, designed around two classic patterns and a zero-dependency Java core.", "size": 12.5, "color": GRAYW, "line": 1.35}])
text(s, SW-M-4.6, 6.10, 4.6, 0.9, [[{"t": NAME1, "size": 13, "color": WHITE, "font": SANS_SB}],
                                   [{"t": NAME2, "size": 13, "color": WHITE, "font": SANS_SB, "space_before": 3}]], align=PP_ALIGN.RIGHT)
notes(s, "BOTH, open together.\n"
    "Thank you. We are " + NAME1 + " and " + NAME2 + ", and this is our SEN3006 project: a Task Management System written in pure Java.\n"
    "The whole system is built around two of the design patterns from the course, Factory Method and Strategy, with a third, State, that the task lifecycle gave us almost for free. There are no external libraries, so the architecture really is the only thing under evaluation today.\n"
    "We want to do four things: define the problem clearly, justify why these patterns fit, walk through the architecture and the running code, and then evaluate the result honestly. " + NAME1 + " starts with the table of contents.")

# 02 TABLE OF CONTENTS --------------------------------------------------------
s = slide(PAPER); eyebrow(s, "Contents"); tracker(s, 1)
title(s, "What this presentation covers.")
toc = [("01","Introduction and problem definition","target-04","indigo",INDIGO,INDIGO_T,"~3 min"),
       ("02","The design patterns: Factory Method, Strategy, State","puzzle-piece-01","violet",VIOLET,VIOLET_T,"~5 min"),
       ("03","Architecture and implementation","grid-01","teal",TEAL,TEAL_T,"~4 min"),
       ("04","Live demonstration","terminal","green",GREEN,GREEN_T,"~3 min"),
       ("05","Results, evaluation and conclusion","trophy-01","amber",AMBER,AMBER_T,"~2 min")]
y0 = 2.62
for i,(n,t,icon,key,acc,tint,tm) in enumerate(toc):
    y = y0 + i*0.74
    rrect(s, M, y, SW-2*M, 0.62, CARD, CARDBD, 1.0, radius=0.06)
    chip(s, M+0.18, y+0.07, 0.48, tint, icon, key)
    text(s, M+0.84, y, 0.7, 0.62, [{"t": n, "size": 14, "color": acc, "font": SERIF_SB}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, M+1.5, y, 8.2, 0.62, [{"t": t, "size": 14.5, "color": INK, "font": SANS_SB}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, SW-M-1.7, y, 1.5, 0.62, [{"t": tm, "size": 11, "color": MUTE, "font": SANS_MD}], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
text(s, M, SH-0.74, SW-2*M, 0.4, [{"t": "Roughly sixteen minutes of talk and demo, then questions. We each take questions on the parts we led.", "size": 10.5, "color": FAINT, "track": 0.4}])
footer(s, 2)
notes(s, NAME1 + ".\n"
    "Here is how the talk is organised, in five sections. Section one defines the problem and the requirements. Section two is the core of the talk and the largest part of the grade: the design patterns, why we chose them, and how we built them. Section three is the architecture and how it maps to Java and SOLID. Section four is a live run of the system. Section five is an honest evaluation and the conclusion.\n"
    "Please hold questions for the end unless something on screen is unclear. Section one is mine.")

# 03 SECTION 01 DIVIDER -------------------------------------------------------
s = slide(DARK)
divider(s, "Section 01", "Introduction and problem definition.",
        "The domain we chose, the objectives we set, and the two problems that drove every pattern choice.",
        "4A45C7", "target-04", 1)
notes(s, NAME1 + ".\n"
    "Section one. Before any pattern or class, we want to pin down the problem, because the assignment is graded partly on a clear problem definition. I will cover the domain and our objectives, the functional and non-functional requirements, and then the two specific problems that the patterns were chosen to solve.")

# 04 INTRODUCTION -------------------------------------------------------------
s = slide(PAPER); eyebrow(s, "Introduction  ·  background and objectives"); tracker(s, 1)
title(s, "Why a task management system?")
sub(s, "Software teams track many kinds of work. How that work gets created and ordered is where a tidy codebase usually starts to sprawl.")
intro=[("clipboard-check","indigo",INDIGO_T,"The domain","A tracker for a software team: bug, feature, and documentation tasks, each moving through one shared lifecycle from open to done."),
       ("zap","indigo",INDIGO_T,"The motivation","Two recurring causes of code sprawl in this kind of system line up cleanly with two patterns from the course."),
       ("target-04","indigo",INDIGO_T,"Our objective","Support new task types and new orderings, and let them be added without editing code that already works."),
       ("cube-01","indigo",INDIGO_T,"In one line","One engine, three ways to run it, sixteen Java files, zero external dependencies.")]
gap=0.5; colw=(SW-2*M-gap)/2; rh=1.74
for i,(icon,key,tint,h,b) in enumerate(intro):
    c,r=i%2,i//2; x=M+c*(colw+gap); y=2.98+r*(rh+0.22)
    card(s, x, y, colw, rh, tint, icon, key, h, b, tsize=14, bsize=11.5)
footer(s, 4)
notes(s, NAME1 + ".\n"
    "Quick background, top left to bottom right. The domain is a task tracker for a software team that handles three kinds of work, bugs, features, and documentation, all on one lifecycle.\n"
    "Our motivation was not to show off patterns. It was that two real, recurring sources of mess in this kind of system happen to match two patterns we studied. So our objective was concrete and testable: add a new task type or a new ordering rule without touching code that already works.\n"
    "And in one line: one engine, three entry points, sixteen files, no libraries. The next slide states the requirements properly.")

# 05 REQUIREMENTS -------------------------------------------------------------
s = slide(PAPER); eyebrow(s, "Problem definition  ·  requirements"); tracker(s, 1)
title(s, "What the system must do.")
gap=0.5; colw=(SW-2*M-gap)/2
listpanel(s, M, 2.55, colw, 3.05, INDIGO_T, "check-square", "indigo", "Functional requirements", [
    "Create bug, feature, and documentation tasks.",
    "Order the list by priority, deadline, or severity.",
    "Move tasks through a validated lifecycle.",
    "Filter by status and summarise the list."], INDIGO, gap=0.58)
listpanel(s, M+colw+gap, 2.55, colw, 3.05, INDIGO_T, "shield-tick", "indigo", "Non-functional requirements", [
    "Extensible: new types or orderings, no edits to old code.",
    "Pure Java 8+, zero external dependencies.",
    "Runs from the terminal: tests, console, and a GUI.",
    "Explainable: the design has to survive a viva."], INDIGO, gap=0.58)
rrect(s, M, 5.86, SW-2*M, 0.74, PANEL, radius=0.06)
text(s, M+0.3, 5.86, SW-2*M-0.6, 0.74, [{"t": "Two requirements drive the whole design: tasks are created in different ways, and the same list is ordered in different ways.", "size": 12.5, "color": INK, "font": SANS_SB, "line": 1.3}], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 5)
notes(s, NAME1 + ".\n"
    "On the left, the functional requirements, the things a user can do: create the three task types, order the list three ways, move tasks through the lifecycle, and filter and summarise.\n"
    "On the right, the non-functional requirements, which is where the architecture lives. It has to be extensible without editing old code, stay pure Java with no dependencies, run from the terminal in three forms, and be explainable, because we have to defend every choice today.\n"
    "The strip at the bottom is the whole talk in one line: different creation, and different ordering. The next two slides show why each one is genuinely hard.")

# 06 PROBLEM 1 ----------------------------------------------------------------
s = slide(PAPER); eyebrow(s, "Problem definition  ·  problem 1"); tracker(s, 1)
title(s, "Tasks have different creation logic.")
sub(s, "Each task type has different fields and different defaults. Building them by hand ties every caller to the concrete classes.")
code(s, M, 2.95, SW-2*M, [
    'new BugTask(title, priority, "MEDIUM", "");          // severity + steps',
    'new FeatureTask(title, priority, 8, 5);              // effort + value',
    'new DocumentationTask(title, priority, "API", "Developers");'], accent=CLAY)
iconcol3(s, [("The flaw","Every entry point imports and names BugTask, FeatureTask, and DocumentationTask directly."),
             ("The cost","Adding a fourth type means hunting down and editing every place that builds a task."),
             ("What we need","We need to make a task from a type identifier, so the caller never names a concrete class.")],
         y=4.95, specs=[(AMBER,AMBER_T,"alert-triangle","amber"),(RED,RED_T,"git-branch-01","red"),(GREEN,GREEN_T,"magic-wand-01","green")])
footer(s, 6)
notes(s, NAME1 + ".\n"
    "Problem one is creation. Look at the three constructors. They are genuinely different: a bug needs a severity and reproduction steps, a feature needs an effort estimate and a business value, a documentation task needs a type and an audience. Each also has its own sensible defaults.\n"
    "The flaw, on the left: without a factory, every entry point has to import and name each concrete class. The cost, in the middle: when we add a fourth type, we have to find and edit every one of those spots. What we need, on the right: a way to make a task from a simple identifier like the string BUG, without the caller knowing which class that is. That is Factory Method, which we reach in section two.")

# 07 PROBLEM 2 ----------------------------------------------------------------
s = slide(PAPER); eyebrow(s, "Problem definition  ·  problem 2"); tracker(s, 1)
title(s, "Ordering depends on the moment.")
sub(s, "The same task list is sorted differently at different times. The list never changes; only the rule does.")
code(s, M, 2.95, SW-2*M, [
    'switch (sortMode) {',
    '    case DEADLINE: /* sort by due date   */ break;',
    '    case SEVERITY: /* bugs first, by rank */ break;',
    '    case URGENT:   /* sort by priority    */ break;',
    '}'], accent=TEAL)
iconcol3(s, [("The flaw","One switch statement inside the manager grows to carry every ordering rule."),
             ("The cost","Each new context edits tested code, which breaks the Open/Closed principle."),
             ("What we need","The ordering should swap at runtime, and new orderings should drop in without touching the manager.")],
         y=5.02, specs=[(AMBER,AMBER_T,"alert-triangle","amber"),(RED,RED_T,"git-branch-01","red"),(GREEN,GREEN_T,"magic-wand-01","green")])
footer(s, 7)
notes(s, NAME1 + ".\n"
    "Problem two is ordering. Three real moments: sprint planning sorts by deadline, incident response sorts by severity with bugs first, and the daily standup sorts by raw priority. The list of tasks is identical in all three; only the rule changes.\n"
    "The naive fix is the switch statement on screen, sitting inside the manager. The flaw is that one method grows without end. The cost is that every new context makes us open and edit code that already works and is already tested, which is a direct Open/Closed violation. What we need is to swap the ordering at runtime and add new orderings as separate units. That is Strategy. Back to me for section two, the patterns.")

# 08 SECTION 02 DIVIDER -------------------------------------------------------
s = slide(DARK)
divider(s, "Section 02", "The design patterns.",
        "Two patterns chosen to match the two problems, plus a third that the task lifecycle gave us at no extra cost.",
        "7C3AED", "puzzle-piece-01", 2)
# pattern colour legend, on a light key-strip so all three brand hues read on the colour field
rrect(s, M, 6.12, 6.98, 0.54, WHITE, radius=0.16)
leg=[("Factory Method","creational","BC5A2B"),("Strategy","behavioural","0E7C72"),("State","behavioural","7C3AED")]
lx=M+0.30
for lab,cat,hexc in leg:
    rrect(s, lx, 6.34, 0.26, 0.10, C(hexc), radius=0.5)
    text(s, lx+0.36, 6.21, 3.0, 0.36, [[{"t": lab+"  ", "size": 11.5, "color": INK, "font": SANS_SB},
                                        {"t": cat, "size": 10, "color": MUTE}]])
    lx += 0.36 + (2.18 if lab!="Strategy" else 1.66)
notes(s, NAME1 + ".\n"
    "Section two is the longest section and carries most of the grade. For each pattern we follow the same shape the report uses: the definition, why it suits this problem, its advantages, and how we actually built it. We start by mapping the two problems to the two patterns, then take Factory Method, then Strategy, then a short bonus on the lifecycle. The colour code at the bottom follows each pattern through the rest of the deck: clay for Factory, teal for Strategy, violet for State.")

# 09 PATTERN SELECTION --------------------------------------------------------
s = slide(PAPER); eyebrow(s, "Pattern selection", VIOLET); tracker(s, 2)
title(s, "Matching patterns to problems.")
sub(s, "We did not start from the patterns. We started from the two problems, and each one maps cleanly to a single pattern.")
rows=[("package","clay",CLAY,CLAY_T,"Problem 1  ·  different creation logic","Factory Method","Creational","Decouples the caller from the concrete task classes."),
      ("switch-horizontal-01","teal",TEAL,TEAL_T,"Problem 2  ·  context-dependent ordering","Strategy","Behavioural","Makes each ordering a swappable object, chosen at runtime."),
      ("refresh-cw-01","violet",VIOLET,VIOLET_T,"Bonus  ·  the task lifecycle","State (in an enum)","Behavioural","Each status declares its own legal transitions.")]
y0=3.0
for i,(icon,key,acc,tint,prob,pat,cat,why) in enumerate(rows):
    y=y0+i*1.18
    rrect(s, M, y, SW-2*M, 1.0, CARD, CARDBD, 1.0, radius=0.05)
    chip(s, M+0.22, y+0.24, 0.52, tint, icon, key)
    text(s, M+1.0, y, 3.7, 1.0, [{"t": prob, "size": 12, "color": BODY, "line": 1.25}], anchor=MSO_ANCHOR.MIDDLE)
    place_icon(s, "arrow-narrow-right", "mute", M+5.05, y+0.5, 0.34)
    text(s, M+5.4, y, 3.2, 1.0, [[{"t": pat, "size": 16, "color": acc, "font": SERIF_SB}],
                                 [{"t": cat+" pattern", "size": 10, "color": MUTE, "space_before": 1, "font": SANS_MD}]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, M+8.55, y, SW-2*M-8.55-0.3, 1.0, [{"t": why, "size": 11.5, "color": BODY, "line": 1.25}], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 9)
notes(s, NAME1 + ".\n"
    "This slide is the bridge, and it sets the colour code for the rest of the deck. Reading each row left to right: problem one, different creation logic, points to Factory Method, a creational pattern in clay, which decouples the caller from the concrete task classes. Problem two, context-dependent ordering, points to Strategy, a behavioural pattern in teal, which makes each ordering a swappable object chosen at runtime.\n"
    "The third row is a bonus we did not have to do. The lifecycle is naturally a state machine, so we encoded State inside a single enum, in violet. We picked this domain precisely because the mappings are natural, not forced. Now the detail, starting with Factory Method.")

# 10 FACTORY METHOD · definition ---------------------------------------------
s = slide(PAPER); eyebrow(s, "Pattern 1  ·  Factory Method", CLAY); tracker(s, 2)
title(s, "Factory Method, and why it fits.")
fm=[("book-open-01","clay",CLAY_T,"Definition","Define a method for making an object, and let each subclass decide which concrete class it produces."),
    ("puzzle-piece-01","clay",CLAY_T,"Why here","Every task type has its own construction rules and defaults. The caller should not have to know them."),
    ("plus-circle","clay",CLAY_T,"Advantage","A new task type is a new class. Code that creates tasks never changes when a type is added."),
    ("file-02","clay",CLAY_T,"Real world","A document app that opens PDF, DOCX, and TXT files through one Open action, without naming a concrete reader.")]
gap=0.5; colw=(SW-2*M-gap)/2; rh=1.74
for i,(icon,key,tint,h,b) in enumerate(fm):
    c,r=i%2,i//2; x=M+c*(colw+gap); y=2.55+r*(rh+0.26)
    card(s, x, y, colw, rh, tint, icon, key, h, b, tsize=14)
footer(s, 10)
notes(s, NAME1 + ".\n"
    "Factory Method, following the report structure. The definition: define a method for making an object, and let each subclass decide which concrete class it returns. We have an abstract TaskFactory with a createTask method and three subclasses that each build their own type.\n"
    "Why it fits here: each task type has different construction rules and defaults, and we do not want the caller carrying that knowledge. The advantage is the one tied to our objective: a new task type is just a new class, and no existing creation code changes. A familiar example is a document editor that opens several formats through one Open action; you do not call a different method per format. Next, how it looks in our code.")

# 11 FACTORY METHOD · code ----------------------------------------------------
s = slide(PAPER); eyebrow(s, "Pattern 1  ·  Factory Method", CLAY); tracker(s, 2)
title(s, "How we built it.")
sub(s, "The manager keeps a registry of factories keyed by a string. The caller asks for a type by name and never sees a concrete class.")
code(s, M, 2.9, SW-2*M, [
    'manager.createTask("BUG", "Login broken", "OAuth callback fails", 8);',
    'manager.createTask("FEATURE", "Dark mode", "Top user request", 5);'], accent=CLAY)
gap=0.5; colw=(SW-2*M-gap)/2
card(s, M, 4.62, colw, 1.78, CLAY_T, "package", "clay", "An abstract class, not an interface",
     "Each subclass owns its defaults: a bug is MEDIUM severity, a feature is 8h of effort at value 5, a doc targets developers.", tsize=13.5, bsize=11.5)
card(s, M+colw+gap, 4.62, colw, 1.78, CLAY_T, "layers-three-01", "indigo", "A template method on the base",
     "createTaskWithDeadline wraps createTask and attaches a deadline once. Written once, reused by all three factories.", tsize=13.5, bsize=11.5)
footer(s, 11)
notes(s, NAME1 + ".\n"
    "Here is the implementation. The manager holds a Map from a string key to a factory. The caller writes createTask with the string BUG, and the manager looks up the right factory and calls it. The caller never imports BugTask. When we added Documentation late in the project, we registered it once and the console menu offered it automatically, because the menu just walks the registry.\n"
    "Two decisions we expect questions on. Left: we used an abstract class, not an interface, so each subclass can own its defaults. Right: createTaskWithDeadline lives on the base and wraps the abstract createTask, which is Template Method sitting on top of Factory Method, so the deadline logic is written once, not three times. Now, Strategy.")

# 12 STRATEGY · definition ----------------------------------------------------
s = slide(PAPER); eyebrow(s, "Pattern 2  ·  Strategy", TEAL); tracker(s, 2)
title(s, "Strategy, and why it fits.")
st=[("book-open-01","teal",TEAL_T,"Definition","Define a family of interchangeable algorithms, wrap each one, and select between them at runtime."),
    ("puzzle-piece-01","teal",TEAL_T,"Why here","The same task list must be ordered differently in different moments, while the app is running."),
    ("plus-circle","teal",TEAL_T,"Advantage","Add or switch an ordering without touching the code that holds and manages the list."),
    ("compass","teal",TEAL_T,"Real world","Like picking a route in a maps app: fastest, shortest, or no tolls, over the same map.")]
gap=0.5; colw=(SW-2*M-gap)/2; rh=1.74
for i,(icon,key,tint,h,b) in enumerate(st):
    c,r=i%2,i//2; x=M+c*(colw+gap); y=2.55+r*(rh+0.26)
    card(s, x, y, colw, rh, tint, icon, key, h, b, tsize=14)
footer(s, 12)
notes(s, NAME1 + ".\n"
    "Strategy, same structure. The definition: define a family of interchangeable algorithms, put each behind a common interface, and choose between them at runtime. Our interface is PriorityStrategy with a single sort method, and we have three implementations.\n"
    "Why it fits: the same list needs different orderings in different moments, and the user changes that while the app runs. The advantage, again tied to our objective: a new ordering is a new class, and the manager that holds the list never changes. The everyday example is a maps app, where fastest, shortest, and avoid tolls are three strategies over one map. Next, the code.")

# 13 STRATEGY · code ----------------------------------------------------------
s = slide(PAPER); eyebrow(s, "Pattern 2  ·  Strategy", TEAL); tracker(s, 2)
title(s, "How we built it.")
sub(s, "Each ordering is a swappable object behind one interface. Changing the sort is a single setter call; the manager never learns which one is active.")
code(s, M, 2.9, SW-2*M, [
    'manager.setPriorityStrategy(new SeverityFirstStrategy());',
    '// the next call to getPrioritizedTasks() returns the new order'], accent=TEAL)
gap=0.5; colw=(SW-2*M-gap)/2
card(s, M, 4.5, colw, 1.9, TEAL_T, "switch-horizontal-01", "teal", "An interface, not a raw Comparator",
     "PriorityStrategy names the role. SeverityFirst also lifts bugs to the top before sorting the rest, which a pairwise Comparator handles awkwardly.", tsize=13.5)
card(s, M+colw+gap, 4.5, colw, 1.9, TEAL_T, "settings-01", "indigo", "A setter, not a constructor argument",
     "Users change the sort mid-session. A setter avoids rebuilding the manager and copying every task into a fresh instance.", tsize=13.5)
footer(s, 13)
notes(s, NAME1 + ".\n"
    "The implementation is small. The manager holds one PriorityStrategy field. Changing the order is one setter call, and the next listing comes back sorted differently, with no recompile.\n"
    "Two decisions. Left: we used a named interface rather than a raw Java Comparator, because the name PriorityStrategy advertises the role in the system, and because SeverityFirst lifts bugs to the top before sorting the rest, which is stateful work a pairwise Comparator does awkwardly. Right: we used a setter rather than constructor injection, because the user swaps sorts while the app runs. The next slide shows this happening for real in the GUI.")

# 14 STRATEGY · live evidence -------------------------------------------------
s = slide(PAPER); eyebrow(s, "Pattern 2  ·  Strategy  ·  evidence", TEAL); tracker(s, 2)
title(s, "The same list, reordered live.")
gap=0.55; colw=(SW-2*M-gap)/2
img_panel(s, "gui-strategy.png", M, 2.5, colw, 3.7, accent=TEAL, caption="Urgent First sort: priority 5 down to 1.")
img_panel(s, "gui-severity.png", M+colw+gap, 2.5, colw, 3.7, accent=TEAL, caption="Severity First sort: both bugs rise to the top.")
footer(s, 14)
notes(s, NAME1 + ".\n"
    "These are two real screenshots of our running GUI, the same five tasks in both. On the left, sorted Urgent First, the rows run by priority, five down to one. On the right, after switching the Sort by dropdown to Severity First, the two bug rows jump to the top, and the status bar confirms the active strategy changed.\n"
    "Nothing was recompiled between these two pictures. The list of tasks is identical; only the strategy object changed, through one setter call wired to the dropdown. This is the clearest single piece of evidence that the pattern works. In the live demo we will do this switch in front of you.")

# 15 BONUS · STATE ------------------------------------------------------------
s = slide(PAPER); eyebrow(s, "Bonus pattern  ·  State", VIOLET); tracker(s, 2)
title(s, "The lifecycle is a state machine.")
img_panel(s, "state-diagram.png", M, 2.5, 5.95, 4.0, accent=VIOLET)
bx=7.15
points=[("TaskStatus owns its transitions","Each enum constant declares its legal next states through canTransitionTo()."),
        ("The base class enforces them","AbstractTask.setStatus() checks that table and throws on an illegal move."),
        ("No shortcuts","The engine refuses to jump OPEN straight to DONE; the cycle has to be walked."),
        ("State, for zero extra files","The whole pattern lives in one enum, with no extra classes.")]
for i,(h,b) in enumerate(points):
    y=2.55+i*1.0
    place_icon(s, "refresh-cw-01", "violet", bx+0.16, y+0.16, 0.32)
    text(s, bx+0.44, y, SW-M-bx-0.44, 0.4, [{"t": h, "size": 13, "color": INK, "font": SANS_SB}])
    text(s, bx+0.44, y+0.30, SW-M-bx-0.44, 0.7, [{"t": b, "size": 11, "color": BODY, "line": 1.28}])
footer(s, 15)
notes(s, NAME1 + ".\n"
    "The bonus. The rubric only asks for one creational and one behavioural pattern, but the lifecycle is a natural state machine, so we encoded it cheaply. The normal path is OPEN, then IN_PROGRESS when a developer picks it up, then REVIEW, then DONE, which is terminal. A reviewer can reject back to IN_PROGRESS, any active task can be BLOCKED, and BLOCKED returns to OPEN.\n"
    "The rules live on the enum itself through canTransitionTo, and setStatus throws if you try an illegal jump such as OPEN straight to DONE. So this is the State pattern in one enum, with no extra classes, and an illegal jump like OPEN to DONE is simply impossible to reach. Over to " + NAME2 + " for the architecture section.")

# 16 SECTION 03 DIVIDER -------------------------------------------------------
s = slide(DARK)
divider(s, "Section 03", "Architecture and implementation.",
        "How the design looks as classes, how one request flows through it, and how it satisfies all five SOLID principles.",
        "0E7C72", "grid-01", 3)
notes(s, NAME2 + ".\n"
    "Section three. We move from the two patterns to the system as a whole. I will show the class diagram and the layers, the sequence of one create-task call, the file structure, and finally how the design lines up with all five SOLID principles.")

# 17 CLASS · FACTORY METHOD ---------------------------------------------------
s = slide(PAPER); eyebrow(s, "Architecture  ·  structure  ·  Factory Method", CLAY); tracker(s, 3)
title(s, "Class view — the Factory Method side.", size=23)
sub(s, "TaskManager talks only to the abstract TaskFactory, which creates a Task. Each concrete factory and each concrete task slots in underneath.", y=2.12, size=11.5, color=MUTE)
img_panel(s, "class-factory.png", M, 2.62, SW-2*M, 4.1, accent=CLAY)
footer(s, 17)
notes(s, NAME2 + ".\n"
    "The full class diagram had fourteen classes side by side, so we split it by pattern. This first half is the Factory Method side. Read it top to bottom: TaskManager depends only on the abstract TaskFactory, never on a concrete factory. Each of the three concrete factories extends TaskFactory, and TaskFactory creates a Task.\n"
    "On the right is the product family the factories build: the Task interface, the AbstractTask base, and the three concrete tasks. The key point is that every arrow points toward an abstraction, the Task interface or the TaskFactory class, which is Dependency Inversion in visual form. That is the structural reason a new task type never touches the manager.")

# 18 CLASS · STRATEGY ---------------------------------------------------------
s = slide(PAPER); eyebrow(s, "Architecture  ·  structure  ·  Strategy", TEAL); tracker(s, 3)
title(s, "Class view — the Strategy side.", size=23)
sub(s, "TaskManager holds one PriorityStrategy and can swap it at runtime. Each ordering is an interchangeable implementation of the same interface.", y=2.12, size=11.5, color=MUTE)
img_panel(s, "class-strategy.png", M, 2.62, SW-2*M, 4.1, accent=TEAL)
footer(s, 18)
notes(s, NAME2 + ".\n"
    "The second half is the Strategy side, and it is deliberately simple. TaskManager holds a single PriorityStrategy reference through its setPriorityStrategy and getPrioritizedTasks methods. It depends only on that interface.\n"
    "Underneath, the three orderings, UrgentFirst, DeadlineFirst, and SeverityFirst, each implement the one sort method. Because the manager only knows the interface, we can swap one ordering for another, or add a fourth, without editing the manager at all. That is the Open/Closed principle falling straight out of the pattern.")

# 19 SEQUENCE · CREATING A TASK -----------------------------------------------
s = slide(PAPER); eyebrow(s, "Architecture  ·  behaviour  ·  creation", CLAY); tracker(s, 3)
title(s, "Creating a task, step by step.", size=23)
sub(s, "The caller passes a type string; the manager looks up the registered factory and returns a Task. No concrete class is ever named by the caller.", y=2.12, size=11.5, color=MUTE)
img_panel(s, "seq-create.png", M, 2.62, SW-2*M, 4.1, accent=CLAY)
footer(s, 19)
notes(s, NAME2 + ".\n"
    "Now the runtime view of the same Factory Method. Trace one createTask call. The entry point, Main, asks the manager for a BUG. The manager looks up the factory registered under that key in its registry, gets the bug factory back, and asks it to build the task. The factory constructs a BugTask and hands it back as a Task reference.\n"
    "Notice what never happens: the caller never writes new BugTask and never imports it. The concrete type stays hidden behind the manager and the factory the whole way through. This is the dynamic counterpart of the class diagram you just saw.")

# 20 SEQUENCE · ORDERING THE LIST ---------------------------------------------
s = slide(PAPER); eyebrow(s, "Architecture  ·  behaviour  ·  ordering", TEAL); tracker(s, 3)
title(s, "Ordering the list, step by step.", size=23)
sub(s, "Set a strategy, then ask for the prioritized list. The manager delegates the sort to whichever strategy is currently set.", y=2.12, size=11.5, color=MUTE)
img_panel(s, "seq-order.png", M, 2.62, SW-2*M, 4.1, accent=TEAL)
footer(s, 20)
notes(s, NAME2 + ".\n"
    "And the runtime view of Strategy. First the caller sets a strategy, here DeadlineFirst, and the manager just stores it as the current strategy. Then the caller asks for the prioritized tasks. The manager does not sort anything itself; it delegates to whichever strategy object is currently set, which sorts by deadline and returns the ordered list.\n"
    "Swap in SeverityFirst and the exact same call returns a different order, with no recompile. That is the whole point of Strategy: the manager owns the list, the strategy owns the rule, and the two vary independently.")

# 21 IMPLEMENTATION OVERVIEW --------------------------------------------------
s = slide(PAPER); eyebrow(s, "Implementation  ·  structure", TEAL); tracker(s, 3)
title(s, "Sixteen files, three layers.")
sub(s, "No package declarations, only java.util and java.time imported. The columns mirror the three layers of the class diagram.")
gap=0.45; colw=(SW-2*M-2*gap)/3
listpanel(s, M, 2.95, colw, 3.3, INDIGO_T, "cube-01", "indigo", "Product layer",
    ["Task (interface)","AbstractTask (base)","BugTask, FeatureTask","DocumentationTask","TaskStatus (enum)"], INDIGO, isize=11.5, gap=0.5)
listpanel(s, M+colw+gap, 2.95, colw, 3.3, VIOLET_T, "puzzle-piece-01", "violet", "Pattern layer",
    ["TaskFactory (abstract)","Bug / Feature / Doc factories","PriorityStrategy (interface)","UrgentFirst, DeadlineFirst","SeverityFirst"], VIOLET, isize=11.5, gap=0.5)
listpanel(s, M+2*(colw+gap), 2.95, colw, 3.3, TEAL_T, "settings-01", "indigo", "Coordination",
    ["TaskManager (coordinator)","Main (automated tests)","TaskManagementApp (console)","gui.TaskManagerGUI (Swing)"], TEAL, isize=11.5, gap=0.5)
rrect(s, M, 6.42, SW-2*M, 0.56, PANEL, radius=0.06)
text(s, M+0.3, 6.42, SW-2*M-0.6, 0.56, [{"t": "All three entry points drive the same engine through its public API.   2 interfaces · 1 enum · 1 abstract · 12 concrete classes.", "size": 11, "color": BODY, "line": 1.2}], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 21)
notes(s, NAME2 + ".\n"
    "The implementation, grouped by the same three layers. The product layer holds the Task interface, the abstract base, the three concrete tasks, and the status enum. The pattern layer holds the abstract factory and its three subclasses, plus the strategy interface and its three implementations. The coordination layer is the single TaskManager, plus the three entry points: Main for automated tests, TaskManagementApp for the console, and the Swing GUI.\n"
    "Sixteen files, two interfaces, one enum, one abstract class, twelve concrete classes, and the only imports are java.util and java.time. Crucially, all three entry points use the same engine through its public API, which is what lets the GUI and the console behave identically.")

# 22 SOLID --------------------------------------------------------------------
s = slide(PAPER); eyebrow(s, "Implementation  ·  principles", TEAL); tracker(s, 3)
title(s, "How the design meets SOLID.")
sol=[("target-01","indigo",INDIGO_T,"S · Single responsibility","Factories build, strategies sort, the manager coordinates. One reason to change each."),
     ("lock-unlocked-01","indigo",INDIGO_T,"O · Open / Closed","A new task type or sort is a new file plus one line. Zero edits to existing classes."),
     ("switch-vertical-01","indigo",INDIGO_T,"L · Liskov substitution","Every factory works through the abstract reference; every strategy through the interface."),
     ("columns-02","indigo",INDIGO_T,"I · Interface segregation","PriorityStrategy declares one method; the Task interface stays minimal."),
     ("anchor","indigo",INDIGO_T,"D · Dependency inversion","The manager's fields are all interfaces and abstract types, never concretes."),
     ("layers-three-01","indigo",INDIGO_T,"The throughline","Abstractions sit in the middle of the system; concrete classes live at the edges.")]
gap=0.45; colw=(SW-2*M-2*gap)/3; rh=1.78
for i,(icon,key,tint,h,b) in enumerate(sol):
    c,r=i%3,i//3; x=M+c*(colw+gap); y=2.5+r*(rh+0.24)
    card(s, x, y, colw, rh, tint, icon, key, h, b, tsize=12.5, bsize=10.8, cs=0.46)
footer(s, 22)
notes(s, NAME2 + ".\n"
    "SOLID, briefly, because it falls out of the two patterns rather than being bolted on. Single responsibility: factories build, strategies sort, the manager coordinates. Open/Closed: a new type or sort is a new file and one line, with zero edits to old code, which is our central objective. Liskov: every factory and strategy stands in for its abstraction without surprises. Interface segregation: the strategy interface has one method, and the Task interface is minimal. Dependency inversion: the manager references only abstractions.\n"
    "The throughline at the bottom sums it up: abstractions in the middle, concrete classes at the edges. If you want one example of inversion, it is that there is no new BugTask anywhere inside the manager. On to the live demonstration.")

# 23 SECTION 04 DIVIDER -------------------------------------------------------
s = slide(DARK)
divider(s, "Section 04", "Live demonstration.",
        "The compiled system, run three ways: the automated test suite, the interactive console, and the Swing GUI.",
        "1E7A45", "terminal", 4)
notes(s, NAME2 + ".\n"
    "Section four is the live demonstration, which the assignment requires. I run the automated tests first, then open the GUI and switch the strategy live, create a task through the factory, and trigger an illegal lifecycle move so you can watch the engine reject it. It is rehearsed to about three minutes, and we have screenshots as a backup if the projector fails.")

# 24 DEMO 1 · TERMINAL --------------------------------------------------------
s = slide(PAPER); eyebrow(s, "Live demonstration  ·  demo 1  ·  the terminal", GREEN); tracker(s, 4)
title(s, "Demo 1 — the terminal.", size=27)
text(s, M, 2.04, SW-2*M, 0.4, [{"t": "Open a terminal in the project root folder, then paste the command for each step.", "size": 12, "color": BODY, "line": 1.3}])
GL=C("D8D4CA"); GG=C("5BD08F"); GM=C("857F73")
rows=[("$ java -cp bin Main", GG),
      ("", GL),
      ("########  Task Management System  ########", GM),
      ("", GL),
      ("  TEST 1   Factory Method Pattern", GL),
      ("     [PASS] correct types, no concretes", GG),
      ("  TEST 2   Strategy Pattern", GL),
      ("     [PASS] one list, three orderings", GG),
      ("  TEST 3   Lifecycle / State machine", GL),
      ("     [PASS] illegal moves are blocked", GG),
      ("  TEST 4-6 edge cases ........  6 / 6", GM),
      ("", GL),
      ("########     ALL TESTS PASSED     ########", GG)]
console(s, M, 2.62, 5.25, 4.06, rows, header="java -cp bin Main")
cx=M+5.62; cw=SW-M-cx
cmdcard(s, cx, 2.62, cw, "STEP 1 · BUILD ONCE", "javac -d bin src/main/java/*.java src/main/java/gui/*.java", GREEN, csize=10.5)
cmdcard(s, cx, 3.66, cw, "STEP 2 · RUN THE AUTOMATED TESTS", "java -cp bin Main", GREEN)
cmdcard(s, cx, 4.70, cw, "STEP 3 · OPEN THE INTERACTIVE CONSOLE", "java -cp bin TaskManagementApp", GREEN)
rrect(s, cx, 5.86, cw, 0.82, PANEL, radius=0.06)
text(s, cx+0.26, 5.86, cw-0.5, 0.82, [[{"t": "Inside the console menu:", "size": 10.5, "color": INK, "font": SANS_SB}],
     [{"t": "1 create a task (Factory)   ·   4 change the sort (Strategy)   ·   3 view sorted   ·   5 change status (State)", "size": 9.5, "color": BODY, "line": 1.25, "space_before": 3}]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 24)
notes(s, NAME2 + ".\n"
    "Demo one is the terminal, because it proves the architecture and not just the interface. Open a terminal in the project root. Step one, paste the javac line to build once. Step two, paste java -cp bin Main: six sections scroll past and it ends on ALL TESTS PASSED, exactly the screen on the left, so every pattern is checked before we touch anything by hand.\n"
    "Step three, paste java -cp bin TaskManagementApp to open the interactive console on the very same engine. From the menu: option 1 creates a task and the Type you enter chooses the factory; option 4 switches the strategy and option 3 re-lists it, so the order changes live with no recompile; option 5 walks a task through its lifecycle. Then on to the GUI.")

# 25 DEMO 2 · GRAPHICAL USER INTERFACE ----------------------------------------
s = slide(PAPER); eyebrow(s, "Live demonstration  ·  demo 2  ·  graphical user interface", GREEN); tracker(s, 4)
title(s, "Demo 2 — the graphical interface.", size=25)
text(s, M, 2.04, SW-2*M, 0.4, [{"t": "From the same project root, one command launches the whole interface — the rest is clicks.", "size": 12, "color": BODY, "line": 1.3}])
img_panel(s, "gui-strategy.png", M, 2.62, 5.25, 4.06, accent=GREEN)
cx=M+5.62; cw=SW-M-cx
cmdcard(s, cx, 2.62, cw, "PASTE THIS TO LAUNCH THE GUI", "java -jar TaskManagerGUI.jar", GREEN, csize=15, h=1.08)
steps=[("browser","indigo",INDIGO_T,"Load the demo data","Demo ▸ Load Strategy Demo seeds the table with tasks."),
       ("switch-horizontal-01","teal",TEAL_T,"Strategy, live","Sort by: Urgent, Deadline, Severity — the table reorders, no recompile."),
       ("package-plus","clay",CLAY_T,"Factory, live","Add a task; the Type field picks which factory builds it."),
       ("refresh-cw-01","violet",VIOLET_T,"State machine","Mark OPEN ▸ DONE: refused. OPEN ▸ IN_PROGRESS: allowed.")]
for i,(icon,key,tint,h,b) in enumerate(steps):
    y=4.06+i*0.72
    chip(s, cx, y, 0.44, tint, icon, key, frac=0.58)
    text(s, cx+0.60, y-0.04, cw-0.60, 0.3, [{"t": h, "size": 12, "color": INK, "font": SANS_SB}])
    text(s, cx+0.60, y+0.23, cw-0.60, 0.42, [{"t": b, "size": 10.2, "color": BODY, "line": 1.18}])
footer(s, 25)
notes(s, NAME2 + ".\n"
    "Demo two is the graphical interface, the same engine behind a Swing window. There is only one thing to paste: java -jar TaskManagerGUI.jar, run from the project root. Then load the demo data from the Demo menu so the table has tasks in it.\n"
    "After that it is all clicks, each tied to a pattern. Cycle the Sort by dropdown through Urgent, Deadline, and Severity and the table reorders instantly, which is Strategy. Add a task from the form, where the Type field decides which factory builds it, which is Factory Method. Finally select a row and try to mark an OPEN task DONE: the engine refuses and explains why, then OPEN to IN_PROGRESS is accepted, which is the State machine. On to the evaluation.")

# 26 SECTION 05 DIVIDER -------------------------------------------------------
s = slide(DARK)
divider(s, "Section 05", "Results, evaluation and conclusion.",
        "What the patterns delivered in practice, the honest limits of this build, and the main lesson.",
        "9A6304", "trophy-01", 5)
notes(s, NAME2 + ".\n"
    "The final section. We measure the result against the objective we set, we are honest about the limitations and the alternatives we considered, and we close with the conclusion and what we learned.")

# 27 RESULTS ------------------------------------------------------------------
s = slide(PAPER); eyebrow(s, "Results  ·  what the patterns delivered", AMBER); tracker(s, 5)
title(s, "New types and orderings, without editing existing code.", size=27)
gap=0.5; colw=(SW-2*M-gap)/2
listpanel(s, M, 2.55, colw, 3.7, GREEN_T, "trend-up-01", "green", "Extension, in numbers", [
    "New task type: 2 new files, 1 line of registration, 0 edits.",
    "New ordering: 1 new file, 1 setter call, 0 edits.",
    "Test 5 in Main asserts this Open/Closed property at runtime."], GREEN, isize=12, gap=0.66)
listpanel(s, M+colw+gap, 2.55, colw, 3.7, GREEN_T, "check-circle", "green", "What worked in practice", [
    "Documentation type was added late, with zero edits to existing code.",
    "The GUI dropdown proves the Strategy swap, live.",
    "The lifecycle enum makes every illegal transition unreachable.",
    "Main doubles as the test harness, with no JUnit needed."], GREEN, isize=12, gap=0.66)
footer(s, 27)
notes(s, NAME2 + ".\n"
    "Did the design meet its objective? On the left, in numbers: a new task type costs two new files and one registration line, with zero edits to existing code; a new ordering costs one new file and one setter call, again zero edits. We did not just claim this; Test 5 in Main checks the property at runtime.\n"
    "On the right, what actually happened during the project. We added the Documentation type late, with no edits to existing code. The GUI dropdown demonstrates the strategy swap live. The lifecycle enum blocks illegal transitions at runtime, before bad state can spread. And Main served as our test harness without pulling in JUnit, which kept the zero-dependency rule.")

# 28 LIMITATIONS --------------------------------------------------------------
s = slide(PAPER); eyebrow(s, "Evaluation  ·  limitations and alternatives", AMBER); tracker(s, 5)
title(s, "An honest look at the design.")
iconcol3(s, [("Limitations","All data is in memory; nothing persists. Single user, no concurrency. Validation throws on the first error rather than collecting them all."),
             ("Alternatives we weighed","Abstract Factory: rejected, there is only one product family. A raw Comparator: loses the named role and the stateful sort."),
             ("Future work","Persistence via Repository, tags via Specification, async dispatch, and a small REST layer over the same engine.")],
         y=2.6, specs=[(AMBER,AMBER_T,"alert-circle","amber"),(INDIGO,INDIGO_T,"intersect-circle","indigo"),(INDIGO,INDIGO_T,"clock-fast-forward","indigo")])
rrect(s, M, 5.5, SW-2*M, 0.92, PANEL, radius=0.06)
text(s, M+0.3, 5.5, SW-2*M-0.6, 0.92, [{"t": "Every limitation above is a scoping choice, not a structural one. The engine is decoupled enough that each item is an addition, not a rewrite.", "size": 12, "color": INK, "font": SANS_MD, "line": 1.35}], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 28)
notes(s, NAME2 + ".\n"
    "An honest evaluation, which the rubric asks for. Limitations: everything is in memory so nothing persists, it is single-user with no concurrency control, and validation throws on the first error rather than collecting them. These are real, and they are deliberate scoping choices for an academic project.\n"
    "Alternatives we weighed: we considered Abstract Factory but rejected it, because we have only one product family, a single Task, so the extra abstraction would not earn its keep. We considered a raw Comparator instead of Strategy, but that loses the named role and handles the bugs-first sort awkwardly. Future work would be persistence with a Repository, tags with a Specification, async dispatch, and a REST layer. Because the engine is decoupled, each of these is an addition, not a rewrite.")

# 29 CONCLUSION ---------------------------------------------------------------
s = slide(PAPER); eyebrow(s, "Conclusion", AMBER); tracker(s, 5)
title(s, "What we set out to do, and what we learned.")
concl=[("check-verified-01","green",GREEN_T,"What we achieved","A working task system where new types and new orderings cost no edits to existing, tested code."),
       ("puzzle-piece-01","violet",VIOLET_T,"The patterns earned their place","Each was chosen to solve a concrete problem from section one, not to satisfy a checklist."),
       ("book-open-01","indigo",INDIGO_T,"The lesson","Naming a role, a factory or a strategy, makes the codebase explain itself to the next reader."),
       ("trophy-01","amber",AMBER_T,"The bigger point","Design patterns are a shared vocabulary that lets a team extend software without breaking what already works.")]
gap=0.5; colw=(SW-2*M-gap)/2; rh=1.74
for i,(icon,key,tint,h,b) in enumerate(concl):
    c,r=i%2,i//2; x=M+c*(colw+gap); y=2.55+r*(rh+0.26)
    card(s, x, y, colw, rh, tint, icon, key, h, b, tsize=13.5)
footer(s, 29)
notes(s, NAME2 + ".\n"
    "To conclude. What we achieved: a working system that meets the objective we set, new types and orderings without editing existing code. The patterns earned their place, because each one answered a specific problem from section one rather than ticking a box.\n"
    "Our main lesson was about naming: calling something a factory or a strategy makes the code explain itself to the next developer who reads it. And the bigger point, the reason the course exists: design patterns are a shared vocabulary that lets a team extend software without breaking what already works. Thank you, and we are happy to take questions.")

# 30 THANK YOU (no names) -----------------------------------------------------
s = slide(DARK)
text(s, M, 2.5, 11.9, 1.8, [[{"t": "Thank you.", "size": 48, "color": WHITE, "font": SERIF_SB, "line": 1.04}],
                            [{"t": "Questions welcome.", "size": 48, "color": INDIGO_D, "font": SERIF_SB, "line": 1.04, "space_before": 4}]])
hline(s, M, 4.95, SW-2*M, RULE_DK, 1.0)
text(s, M, 5.18, 9.8, 1.0, [{"t": "We can speak to any part: the problem, the patterns, the architecture, the code, or the demo.", "size": 13, "color": GRAYW, "line": 1.4}])
# pattern legend echo
leg=[("Factory Method", CLAY),("Strategy", TEAL),("State", VIOLET)]
lx=M
for lab,col in leg:
    rrect(s, lx, 6.18, 0.26, 0.10, col, radius=0.5)
    text(s, lx+0.36, 6.06, 2.6, 0.34, [{"t": lab, "size": 11, "color": GRAYW, "font": SANS_SB}])
    lx += 0.36 + (2.0 if lab!="Strategy" else 1.5)
text(s, SW-M-3.4, 6.06, 3.4, 0.34, [{"t": "SEN3006 · SOFTWARE ARCHITECTURE", "size": 9.5, "color": GRAYW, "track": 1.4, "font": SANS_SB}], align=PP_ALIGN.RIGHT)
notes(s, "Both.\n"
    "Thank you. We will split questions by area: " + NAME1 + " takes the problem framing and the three patterns, Factory Method, Strategy, and State; " + NAME2 + " takes the architecture, the live demo, SOLID, and the evaluation.\n"
    "Anticipated questions. Why an abstract factory class instead of an interface? Per-subclass defaults plus Template Method. Why a strategy interface instead of a Comparator? It advertises the role and supports the partition-then-sort. Why a string-keyed registry? Runtime lookup from any entry point. Why no JUnit? The zero-dependency rule, with Main as the harness. The study guide in docs/design has a fuller cheat sheet.")

# 31 APPENDIX · component view -----------------------------------------------
s = slide(PAPER); eyebrow(s, "Appendix  ·  supporting view"); tracker(s, 5)
title(s, "Component view.", size=27)
sub(s, "The module structure: the UI depends on the manager, which delegates to the factory and strategy modules over the shared domain model.")
img_panel(s, "component-diagram.png", M, 3.0, SW-2*M, 3.0, accent=INDIGO)
text(s, M, 6.18, SW-2*M, 0.3, [{"t": "Each box is a module; every arrow is a compile-time dependency pointing inward toward the manager and the abstractions.", "size": 9.5, "color": FAINT, "track": 0.4}], align=PP_ALIGN.CENTER)
footer(s, 31)
notes(s, "Appendix, for questions.\n"
    "A supporting view kept in reserve. The component diagram groups the system into the UI, manager, factory, strategy, and domain modules. Every dependency points inward toward the manager and the abstractions, which is the module-level echo of the class diagram. It is explained in full in the report.")

# 32 APPENDIX · deployment view ----------------------------------------------
s = slide(PAPER); eyebrow(s, "Appendix  ·  supporting view"); tracker(s, 5)
title(s, "Deployment view.", size=27)
sub(s, "One JVM, one shared engine, three entry points. No network, no database, no external services.")
img_panel(s, "deployment-diagram.png", 2.7, 2.95, SW-2*2.7, 3.5, accent=INDIGO)
footer(s, 32)
notes(s, "Appendix, for questions.\n"
    "The last supporting view. The deployment diagram shows the whole system runs in a single Java virtual machine, with the three entry points, Main, the console app, and the GUI, all driving one shared engine. There is no network and no external service, which is what zero dependencies means in practice. It is explained in full in the report.")

prs.save(OUT); print("SAVED", OUT, "|", len(prs.slides._sldIdLst), "slides")
